"""
Build ClearLaunch UVP Summary Deck (.pptx) — 9 slides
Matches the ICP Summary Deck styling pattern exactly.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Brand colors
GREEN = RGBColor(0x1B, 0x9B, 0x5E)
BRIGHT = RGBColor(0x3B, 0xEB, 0x96)
DARK = RGBColor(0x12, 0x17, 0x18)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF6, 0xF3, 0xEF)
LGGREEN = RGBColor(0xE5, 0xF5, 0xEC)
ALTROW = RGBColor(0xF0, 0xF9, 0xF4)
BORDER = RGBColor(0xA8, 0xD9, 0xBD)

# Slide dimensions (matching ICP deck)
SLIDE_W = 9144000
SLIDE_H = 5143500

# Common measurements
MARGIN_L = Emu(548640)
CONTENT_W = Emu(8046720)
HEADER_H = Emu(777240)
FOOTER_Y = Emu(4873752)
FOOTER_H = Emu(269748)


def add_header_banner(slide, title_text):
    """Green header banner with white title text."""
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    banner.fill.solid()
    banner.fill.fore_color.rgb = GREEN
    banner.line.fill.background()

    title = slide.shapes.add_textbox(MARGIN_L, Emu(137160), CONTENT_W, Emu(502920))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Calibri"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide, subtitle="UVP Summary"):
    """Dark footer bar with ClearThink branding."""
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, FOOTER_Y, SLIDE_W, FOOTER_H)
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = DARK
    footer_bar.line.fill.background()

    footer_text = slide.shapes.add_textbox(MARGIN_L, Emu(4892040), Emu(8229600), Emu(228600))
    tf = footer_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"ClearThink Marketing  |  {subtitle}"
    p.font.name = "Calibri"
    p.font.size = Pt(8)
    p.font.color.rgb = LGGREEN


def add_content_card(slide, left, top, width, height, label, placeholder,
                     accent_color=GREEN):
    """Rounded rectangle card with accent bar, label, and placeholder text."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)

    # Left accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Emu(top + 25400), Emu(50800),
                                     Emu(height - 50800))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    # Label
    label_box = slide.shapes.add_textbox(Emu(left + 190500), Emu(top + 101600),
                                          Emu(width - 254000), Emu(228600))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left + 190500),
                                      Emu(top + 330200), Emu(width - 381000), Emu(19050))
    divider.fill.solid()
    divider.fill.fore_color.rgb = LGGREEN
    divider.line.fill.background()

    # Placeholder text
    ph_box = slide.shapes.add_textbox(Emu(left + 190500), Emu(top + 406400),
                                       Emu(width - 381000), Emu(height - 508000))
    tf2 = ph_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = placeholder
    p2.font.name = "Calibri"
    p2.font.size = Pt(10)
    p2.font.color.rgb = DARK


def add_numbered_badge(slide, left, top, number, bg_color=GREEN):
    """Small numbered badge (like tier badges)."""
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(381000), Emu(254000))
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()

    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(number)
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE


# ========================================
# BUILD THE DECK
# ========================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank

# ========================================
# SLIDE 1: Title
# ========================================
slide = prs.slides.add_slide(blank_layout)

# Thin top bar
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(50800))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = GREEN
top_bar.line.fill.background()

# Title
title = slide.shapes.add_textbox(MARGIN_L, Emu(1005840), CONTENT_W, Emu(731520))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "UNIQUE VALUE PROPOSITION"
p.font.name = "Calibri"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK

# UVP badge
badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(1828800),
                                Emu(640080), Emu(274320))
badge.fill.solid()
badge.fill.fore_color.rgb = GREEN
badge.line.fill.background()
badge_text = slide.shapes.add_textbox(MARGIN_L, Emu(1828800), Emu(640080), Emu(274320))
tf = badge_text.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "UVP"
p.font.name = "Calibri"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE

# Client name
client = slide.shapes.add_textbox(MARGIN_L, Emu(2286000), CONTENT_W, Emu(457200))
tf = client.text_frame
p = tf.paragraphs[0]
p.text = "[Client Name]"
p.font.name = "Calibri"
p.font.size = Pt(20)
p.font.color.rgb = DARK

# Divider
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(2926080),
                                  Emu(2286000), Emu(25400))
divider.fill.solid()
divider.fill.fore_color.rgb = BRIGHT
divider.line.fill.background()

# Subtitle
sub = slide.shapes.add_textbox(MARGIN_L, Emu(3200400), CONTENT_W, Emu(548640))
tf = sub.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ClearLaunch System  |  Step 4: UVP Development\n[Date]"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = DARK

add_footer(slide, "UVP Summary")

# ========================================
# SLIDE 2: Problem & Market Gap
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "PROBLEM & MARKET GAP")

card_top = 1051560
card_h = 1143000
card_gap = 152400

add_content_card(slide, MARGIN_L, card_top, Emu(3931920), card_h,
                 "Core Problem",
                 "[What specific problem does this business solve that others don't address effectively?]")

add_content_card(slide, Emu(4663440), card_top, Emu(3931920), card_h,
                 "Market Gaps",
                 "[What gaps or underserved needs does this business specifically address?]")

add_content_card(slide, MARGIN_L, card_top + card_h + card_gap,
                 CONTENT_W, Emu(1600200),
                 "Pain Points Eliminated",
                 "[What specific pain points do they eliminate for customers that others typically overlook?]")

add_footer(slide)

# ========================================
# SLIDE 3: Methodology & Approach
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "METHODOLOGY & APPROACH")

add_content_card(slide, MARGIN_L, card_top, CONTENT_W, card_h,
                 "Unique Methodology",
                 "[What unique methodology, process, or approach sets their offerings apart from competitors?]")

add_content_card(slide, MARGIN_L, card_top + card_h + card_gap,
                 Emu(3931920), Emu(1600200),
                 "Delivery Difference",
                 "[How do they deliver differently than competitors? (faster, more transparent, more automated, etc.)]")

add_content_card(slide, Emu(4663440), card_top + card_h + card_gap,
                 Emu(3931920), Emu(1600200),
                 "Comprehensive Solution",
                 "[What unique combination of products/services provides a comprehensive solution?]")

add_footer(slide)

# ========================================
# SLIDE 4: Expertise & Credibility
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "EXPERTISE & CREDIBILITY")

add_content_card(slide, MARGIN_L, card_top, CONTENT_W, Emu(1524000),
                 "Specialized Expertise",
                 "[What specialized expertise, certifications, technology, or experience does the team possess that's rare in their space?]")

add_content_card(slide, MARGIN_L, card_top + 1524000 + card_gap,
                 CONTENT_W, Emu(1524000),
                 "Culture & Values → Customer Benefits",
                 "[What aspects of company culture or values translate into tangible benefits for customers?]")

add_footer(slide)

# ========================================
# SLIDE 5: Results & Outcomes
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "RESULTS & OUTCOMES")

add_content_card(slide, MARGIN_L, card_top, CONTENT_W, Emu(1524000),
                 "Unique Results",
                 "[What specific results or outcomes can customers expect that they can't get elsewhere?]")

add_content_card(slide, MARGIN_L, card_top + 1524000 + card_gap,
                 CONTENT_W, Emu(1524000),
                 "Metrics & KPIs Improved",
                 "[What specific metrics or KPIs can they improve for customers better than competitors?]")

add_footer(slide)

# ========================================
# SLIDE 6: Risk Reduction & Assurance
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "RISK REDUCTION & ASSURANCE")

add_content_card(slide, MARGIN_L, card_top, Emu(3931920), Emu(2743200),
                 "Guarantees & Assurances",
                 "[What guarantees, assurances, or risk-reduction elements can they offer that competitors don't?]")

add_content_card(slide, Emu(4663440), card_top, Emu(3931920), Emu(2743200),
                 "Pricing & Value Model",
                 "[How does their pricing model or structure provide better value than traditional approaches?]")

add_footer(slide)

# ========================================
# SLIDE 7: Proof & Validation
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "PROOF & VALIDATION")

add_content_card(slide, MARGIN_L, card_top, CONTENT_W, Emu(1524000),
                 "Customer Feedback & Testimonials",
                 "[What specific customer feedback, testimonials, or reviews highlight their unique advantages?]")

add_content_card(slide, MARGIN_L, card_top + 1524000 + card_gap,
                 CONTENT_W, Emu(1524000),
                 "Case Studies & Evidence",
                 "[What case studies, examples, or data points demonstrate their value proposition in action?]")

add_footer(slide)

# ========================================
# SLIDE 8: Top 3 Differentiators
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "TOP 3 DIFFERENTIATORS")

col_w = Emu(2560320)
col_gap = Emu(182880)

for i in range(3):
    x = int(MARGIN_L) + i * (int(col_w) + int(col_gap))
    y = card_top

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Emu(2926080))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)

    # Number badge
    add_numbered_badge(slide, x + 127000, y + 127000, i + 1)

    # Differentiator name
    name_box = slide.shapes.add_textbox(x + 127000, y + 482600, int(col_w) - 254000, 330200)
    tf = name_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[Differentiator {i+1} Name]"
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + 127000, y + 812800,
                                  int(col_w) - 254000, 19050)
    div.fill.solid()
    div.fill.fore_color.rgb = LGGREEN
    div.line.fill.background()

    # Description
    desc_box = slide.shapes.add_textbox(x + 127000, y + 914400,
                                         int(col_w) - 254000, Emu(1828800))
    tf2 = desc_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "[2-3 sentence explanation of this differentiator, connected to a specific customer pain point]"
    p2.font.name = "Calibri"
    p2.font.size = Pt(9)
    p2.font.color.rgb = DARK

add_footer(slide)

# ========================================
# SLIDE 9: UVP Statement + Elevator Pitch
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "UVP STATEMENT & ELEVATOR PITCH")

# UVP Statement card
card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, card_top,
                                CONTENT_W, Emu(1371600))
card1.fill.solid()
card1.fill.fore_color.rgb = LGGREEN
card1.line.color.rgb = GREEN
card1.line.width = Pt(1)

# UVP accent bar
accent1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, card_top + 25400,
                                  Emu(50800), Emu(1320800))
accent1.fill.solid()
accent1.fill.fore_color.rgb = GREEN
accent1.line.fill.background()

# UVP label
lbl1 = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), card_top + 76200,
                                  Emu(7620000), Emu(228600))
tf = lbl1.text_frame
p = tf.paragraphs[0]
p.text = "Draft UVP Statement"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GREEN

# UVP text
uvp_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), card_top + 355600,
                                     Emu(7620000), Emu(914400))
tf = uvp_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[We help [ICP] achieve [primary outcome] through [unique methodology]. Unlike [competitors] — who [what they do wrong] — we [what you do differently] because [proof/credibility].]"
p.font.name = "Calibri"
p.font.size = Pt(12)
p.font.color.rgb = DARK

# Elevator Pitch card
ep_top = card_top + 1371600 + 152400
card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, ep_top,
                                CONTENT_W, Emu(1371600))
card2.fill.solid()
card2.fill.fore_color.rgb = WHITE
card2.line.color.rgb = BORDER
card2.line.width = Pt(0.5)

accent2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, ep_top + 25400,
                                  Emu(50800), Emu(1320800))
accent2.fill.solid()
accent2.fill.fore_color.rgb = GREEN
accent2.line.fill.background()

lbl2 = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), ep_top + 76200,
                                  Emu(7620000), Emu(228600))
tf = lbl2.text_frame
p = tf.paragraphs[0]
p.text = "Elevator Pitch (30 Seconds)"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GREEN

ep_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), ep_top + 355600,
                                    Emu(7620000), Emu(914400))
tf = ep_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[Natural-sounding paragraph following: name the problem → name the audience → explain what you do differently → prove it with a specific result]"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = DARK

add_footer(slide)

# ========================================
# SAVE
# ========================================
prs.save("Frameworks/ClearLaunch_UVP_Summary_Deck.pptx")
print("Created: Frameworks/ClearLaunch_UVP_Summary_Deck.pptx (9 slides)")
