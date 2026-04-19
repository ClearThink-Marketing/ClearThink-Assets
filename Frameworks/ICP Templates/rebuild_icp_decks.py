"""
Rebuild ClearLaunch ICP Summary Decks (B2B + B2C)
Clean layout with ClearThink brand colors applied.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
GREEN = RGBColor(0x1B, 0x9B, 0x5E)
BRIGHT = RGBColor(0x3B, 0xEB, 0x96)
DARK = RGBColor(0x12, 0x17, 0x18)
CREAM = RGBColor(0xF6, 0xF3, 0xEF)
LGGREEN = RGBColor(0xE5, 0xF5, 0xEC)
ALTROW = RGBColor(0xF0, 0xF9, 0xF4)
BORDER = RGBColor(0xA8, 0xD9, 0xBD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Dimensions ──
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - 2 * MARGIN


def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Helper Functions ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0.75)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.rotation = 0
    # Minimal corner rounding
    shape.adjustments[0] = 0.02
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = border_width
    else:
        ln.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = Pt(0.5)
    else:
        ln.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=10, bold=False,
             color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_divider(slide, left, top, width, color=BORDER, thickness=Pt(1.5)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_header_bar(slide, text):
    """Green header bar across the top of a slide."""
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), fill_color=GREEN)
    add_text(slide, MARGIN, Inches(0.15), CONTENT_W, Inches(0.55), text,
             font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)


def add_footer(slide):
    """Branded footer bar."""
    bar = add_rect(slide, Inches(0), Inches(5.33), SLIDE_W, Inches(0.295), fill_color=DARK)
    add_text(slide, MARGIN, Inches(5.35), Inches(9), Inches(0.25),
             "ClearThink Marketing  |  ICP Summary", font_size=8, color=LGGREEN)


def add_card(slide, left, top, width, height, accent_color=GREEN):
    """A content card with cream background, green left accent, subtle border."""
    card = add_shape(slide, left, top, width, height, fill_color=CREAM, border_color=BORDER)
    # Left accent stripe
    accent = add_rect(slide, left, top + Pt(2), Pt(4), height - Pt(4), fill_color=accent_color)
    return card


def add_label_value_row(slide, left, top, label_w, value_w, row_h, label, value, is_alt=False):
    """A label-value pair row, optionally with alternating background."""
    if is_alt:
        add_rect(slide, left, top, label_w + value_w, row_h, fill_color=ALTROW)
    add_text(slide, left + Pt(6), top + Pt(3), label_w - Pt(12), row_h - Pt(6),
             label, font_size=9, bold=True, color=GREEN)
    add_text(slide, left + label_w + Pt(4), top + Pt(3), value_w - Pt(10), row_h - Pt(6),
             value, font_size=9, color=DARK)


def build_data_card(slide, left, top, width, height, icon_label, section_title, rows, accent_color=GREEN):
    """
    Build a card with title, divider, and label-value rows.
    rows: list of (label, value) tuples
    """
    card = add_card(slide, left, top, width, height, accent_color)

    inner_left = left + Pt(12)
    inner_w = width - Pt(24)
    label_w = Inches(1.5)
    value_w = inner_w - label_w

    # Section title
    add_text(slide, inner_left, top + Pt(8), inner_w, Pt(22),
             section_title, font_size=11, bold=True, color=GREEN)

    # Divider
    add_divider(slide, inner_left, top + Pt(32), inner_w, color=BORDER, thickness=Pt(1))

    # Rows
    row_h = Pt(38)
    y = top + Pt(38)
    for i, (label, value) in enumerate(rows):
        add_label_value_row(slide, inner_left, y, label_w, value_w, row_h, label, value, is_alt=(i % 2 == 1))
        y += row_h


def build_tier_card(slide, left, top, width, height, tier_label, tier_name, description, assessment, accent_color):
    """Tier overview card for slide 2."""
    card = add_shape(slide, left, top, width, height, fill_color=CREAM, border_color=BORDER)

    # Tier badge
    badge_w = Inches(0.7)
    badge_h = Pt(20)
    badge = add_rect(slide, left + Pt(10), top + Pt(10), badge_w, badge_h, fill_color=accent_color)
    add_text(slide, left + Pt(10), top + Pt(10), badge_w, badge_h,
             tier_label, font_size=8, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Name
    add_text(slide, left + Pt(10), top + Pt(38), width - Pt(20), Pt(22),
             tier_name, font_size=12, bold=True, color=DARK)

    # Divider
    add_divider(slide, left + Pt(10), top + Pt(64), width - Pt(20), color=BORDER)

    # Description label + text
    add_text(slide, left + Pt(10), top + Pt(72), width - Pt(20), Pt(16),
             "Description", font_size=8, bold=True, color=GREEN)
    add_text(slide, left + Pt(10), top + Pt(90), width - Pt(20), Pt(45),
             description, font_size=9, color=DARK)

    # Assessment label + text
    add_text(slide, left + Pt(10), top + Pt(142), width - Pt(20), Pt(16),
             "Assessment", font_size=8, bold=True, color=GREEN)
    add_text(slide, left + Pt(10), top + Pt(160), width - Pt(20), Pt(45),
             assessment, font_size=9, color=DARK)


# ═══════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════

def slide_01_title(prs, variant):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, CREAM)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=GREEN)

    # Title
    add_text(slide, MARGIN, Inches(1.1), CONTENT_W, Inches(0.8),
             "IDEAL CUSTOMER PROFILE", font_size=32, bold=True, color=DARK)

    # Variant badge
    badge = add_rect(slide, MARGIN, Inches(2.0), Inches(0.7), Inches(0.3), fill_color=GREEN)
    add_text(slide, MARGIN, Inches(2.0), Inches(0.7), Inches(0.3),
             variant, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Client name
    add_text(slide, MARGIN, Inches(2.5), CONTENT_W, Inches(0.5),
             "[Client Name]", font_size=20, color=DARK)

    # Divider
    add_divider(slide, MARGIN, Inches(3.2), Inches(2.5), color=GREEN, thickness=Pt(2))

    # Subtitle
    add_text(slide, MARGIN, Inches(3.5), CONTENT_W, Inches(0.6),
             "ClearLaunch System  |  Stage 1: ICP & Market Definition\n[Date]",
             font_size=11, color=DARK)

    # Footer
    add_rect(slide, Inches(0), Inches(5.33), SLIDE_W, Inches(0.295), fill_color=DARK)
    add_text(slide, MARGIN, Inches(5.35), Inches(9), Inches(0.25),
             "ClearThink Marketing", font_size=8, color=LGGREEN)


def slide_02_tiers(prs, variant):
    """Customer Types / Segments overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    title = "YOUR CUSTOMER TYPES" if variant == "B2B" else "YOUR CUSTOMER SEGMENTS"
    add_header_bar(slide, title)

    card_w = Inches(2.8)
    card_h = Inches(3.2)
    gap = Inches(0.2)
    start_x = MARGIN
    top = Inches(1.2)

    tier1_name = "TYPE A (PRIMARY)" if variant == "B2B" else "SEGMENT A (PRIMARY)"
    tier2_name = "TYPE B" if variant == "B2B" else "SEGMENT B"
    tier3_name = "TYPE C" if variant == "B2B" else "SEGMENT C"

    t1_desc = "[Best, most profitable customer type]" if variant == "B2B" else "[Highest-value consumer segment]"
    t2_desc = "[Another customer type you serve]" if variant == "B2B" else "[Secondary consumer segment]"
    t3_desc = "[Other customer types]" if variant == "B2B" else "[Other consumer segments]"

    t1_assess = "[Why they're your best]" if variant == "B2B" else "[Why they're your best customers]"
    t2_assess = "[How they compare to Type A]" if variant == "B2B" else "[How they compare to Segment A]"
    t3_assess = "[Worth pursuing or not?]"

    build_tier_card(slide, start_x, top, card_w, card_h, "Tier 1", tier1_name, t1_desc, t1_assess, GREEN)
    build_tier_card(slide, start_x + card_w + gap, top, card_w, card_h, "Tier 2", tier2_name, t2_desc, t2_assess, BRIGHT)
    build_tier_card(slide, start_x + 2 * (card_w + gap), top, card_w, card_h, "Tier 3", tier3_name, t3_desc, t3_assess, BORDER)

    add_footer(slide)


def slide_03_deep_dive(prs, variant):
    """Tier 1 Deep Dive."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    title = "TIER 1 DEEP DIVE: COMPANY PROFILE" if variant == "B2B" else "TIER 1 DEEP DIVE: CUSTOMER PROFILE"
    add_header_bar(slide, title)

    # Profile card
    profile_card = add_card(slide, MARGIN, Inches(1.15), CONTENT_W, Inches(0.9))
    name_placeholder = '[ICP Name \u2014 e.g., "Mid-Market SaaS Company"]' if variant == "B2B" else '[Segment Name \u2014 e.g., "Health-Conscious Millennial Parent"]'
    add_text(slide, MARGIN + Pt(15), Inches(1.25), CONTENT_W - Pt(30), Pt(26),
             name_placeholder, font_size=14, bold=True, color=DARK)
    add_text(slide, MARGIN + Pt(15), Inches(1.6), CONTENT_W - Pt(30), Pt(22),
             "[One-line description of who they are]", font_size=10, color=DARK)

    # Two cards side by side
    card_w = Inches(4.3)
    gap = Inches(0.2)
    card_h = Inches(2.2)
    y = Inches(2.35)

    # Left: Why best
    left_card = add_card(slide, MARGIN, y, card_w, card_h)
    add_text(slide, MARGIN + Pt(15), y + Pt(10), card_w - Pt(30), Pt(22),
             "Why They're Your Best Customer", font_size=11, bold=True, color=GREEN)
    add_divider(slide, MARGIN + Pt(15), y + Pt(36), card_w - Pt(30))
    why_text = "[What makes this type of company ideal for your business?]" if variant == "B2B" else "[What makes this segment ideal \u2014 LTV, referral rate, retention, etc.]"
    add_text(slide, MARGIN + Pt(15), y + Pt(44), card_w - Pt(30), Inches(1.2),
             why_text, font_size=10, color=DARK)

    # Right: Strategic Priority
    right_card = add_card(slide, MARGIN + card_w + gap, y, card_w, card_h)
    add_text(slide, MARGIN + card_w + gap + Pt(15), y + Pt(10), card_w - Pt(30), Pt(22),
             "Strategic Priority", font_size=11, bold=True, color=GREEN)
    add_divider(slide, MARGIN + card_w + gap + Pt(15), y + Pt(36), card_w - Pt(30))

    # Priority badge
    badge = add_rect(slide, MARGIN + card_w + gap + Pt(15), y + Pt(48), Inches(1), Pt(22), fill_color=GREEN)
    add_text(slide, MARGIN + card_w + gap + Pt(15), y + Pt(48), Inches(1), Pt(22),
             "PRIMARY", font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    priority_text = "[Max investment. Direct outreach, personalized campaigns.]" if variant == "B2B" else "[Max investment. Targeted ads, influencer partnerships, community building.]"
    add_text(slide, MARGIN + card_w + gap + Pt(15), y + Pt(78), card_w - Pt(30), Inches(0.8),
             priority_text, font_size=10, color=DARK)

    add_footer(slide)


def slide_04_profile_details(prs, variant):
    """Firmographics (B2B) or Demographics & Psychographics (B2C)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    if variant == "B2B":
        add_header_bar(slide, "FIRMOGRAPHICS")
        left_title = "Organization"
        left_rows = [
            ("Industry / Vertical", "[e.g., Construction, SaaS]"),
            ("Company Size", "[e.g., 10-50 employees]"),
            ("Annual Revenue", "[e.g., $1M-$5M]"),
            ("Years in Business", "[Range]"),
            ("Growth Stage", "[Startup / Scaling / Established]"),
        ]
        right_title = "Market Position"
        right_rows = [
            ("Geographic Markets", "[Cities, states, regions]"),
            ("Digital Maturity", "[Minimal / Some / Advanced]"),
            ("Org Structure", "[Owner-operated / Departmentalized]"),
            ("Budget Range", "[Per project or annual spend]"),
            ("Budget Cycle", "[When they plan & allocate]"),
        ]
    else:
        add_header_bar(slide, "DEMOGRAPHICS & PSYCHOGRAPHICS")
        left_title = "Demographics"
        left_rows = [
            ("Age Range", "[e.g., 28-42]"),
            ("Gender", "[If relevant to product]"),
            ("Household Income", "[e.g., $75K-$150K]"),
            ("Location", "[Urban / Suburban / Regional]"),
            ("Life Stage", "[Single / Married / Parent / Empty Nest]"),
        ]
        right_title = "Psychographics"
        right_rows = [
            ("Values", "[What matters most to them]"),
            ("Lifestyle", "[How they spend time & money]"),
            ("Identity", "[How they see themselves]"),
            ("Influences", "[Who/what shapes their decisions]"),
            ("Digital Behavior", "[Social platforms, content habits]"),
        ]

    card_w = Inches(4.3)
    card_h = Inches(3.8)
    gap = Inches(0.2)
    y = Inches(1.1)

    build_data_card(slide, MARGIN, y, card_w, card_h, "", left_title, left_rows)
    build_data_card(slide, MARGIN + card_w + gap, y, card_w, card_h, "", right_title, right_rows)

    add_footer(slide)


def slide_05_pain_triggers(prs, variant):
    """Needs, Pain Points & Buying/Purchase Triggers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    if variant == "B2B":
        add_header_bar(slide, "NEEDS, PAIN POINTS & BUYING TRIGGERS")
        left_title = "Pain Points"
        left_rows = [
            ("Primary Need", "[Core problem they hire you to solve]"),
            ("Secondary Needs", "[Additional problems or goals]"),
            ("Current Solution", "[What they do now]"),
            ("Pain with Current", "[Why it isn't working]"),
            ("Urgency Level", "[Critical / Important / Nice-to-Have]"),
        ]
        right_title = "Buying Triggers"
        right_rows = [
            ("Triggers", "[New funding, growth, bad vendor, etc.]"),
            ("Decision Timeline", "[First contact to signed deal]"),
            ("Decision Maker", "[Title/role of the signer]"),
            ("Influencers", "[Who else has input?]"),
            ("Key Objections", "[Common pushback in sales]"),
        ]
    else:
        add_header_bar(slide, "NEEDS, PAIN POINTS & PURCHASE TRIGGERS")
        left_title = "Pain Points"
        left_rows = [
            ("Primary Need", "[Core problem your product solves]"),
            ("Emotional Driver", "[How the problem makes them feel]"),
            ("Current Solution", "[What they use now / do without]"),
            ("Frustration", "[Why current option falls short]"),
            ("Urgency Level", "[Immediate / Seasonal / Ongoing]"),
        ]
        right_title = "Purchase Triggers"
        right_rows = [
            ("Trigger Events", "[Life change, season, recommendation]"),
            ("Decision Speed", "[Impulse / Considered / Researched]"),
            ("Price Sensitivity", "[Budget-driven / Value-driven / Premium]"),
            ("Influencer Role", "[Spouse, friends, online reviews]"),
            ("Key Objections", "[Common reasons they hesitate]"),
        ]

    card_w = Inches(4.3)
    card_h = Inches(3.8)
    gap = Inches(0.2)
    y = Inches(1.1)

    build_data_card(slide, MARGIN, y, card_w, card_h, "", left_title, left_rows)
    build_data_card(slide, MARGIN + card_w + gap, y, card_w, card_h, "", right_title, right_rows)

    add_footer(slide)


def slide_06_landscape(prs, variant):
    """Beachhead & Competitive (B2B) or Customer Journey & Touchpoints (B2C)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    if variant == "B2B":
        add_header_bar(slide, "BEACHHEAD & COMPETITIVE LANDSCAPE")
        left_title = "Beachhead Selection"
        left_rows = [
            ("Beachhead Segment", "[Specific niche to dominate first]"),
            ("Why Start Here", "[Why this segment first?]"),
            ("Pain Acute Enough?", "[Will they pay to solve now?]"),
            ("Niche Winnable?", "[No dominant competitor?]"),
            ("Expansion Path", "[Adjacent segments after winning]"),
        ]
        right_title = "Competitive Landscape"
        right_rows = [
            ("Direct Competitors", "[Who else serves this segment?]"),
            ("Market Alternative", "[Incumbent controlling the budget]"),
            ("Product Alternative", "[Innovative solution from other category]"),
            ("Your Differentiation", "[Why would they pick you?]"),
            ("Switching Barriers", "[What keeps them with current?]"),
        ]
    else:
        add_header_bar(slide, "CUSTOMER JOURNEY & TOUCHPOINTS")
        left_title = "Discovery & Research"
        left_rows = [
            ("Awareness", "[How they first hear about you]"),
            ("Research", "[Where they look for information]"),
            ("Comparison", "[How they evaluate alternatives]"),
            ("Trust Builders", "[Reviews, testimonials, social proof]"),
            ("Content That Works", "[Video, blog, social, UGC]"),
        ]
        right_title = "Conversion & Retention"
        right_rows = [
            ("Purchase Channel", "[Online / In-store / App / DTC]"),
            ("First Purchase", "[What they buy first]"),
            ("Avg. Order Value", "[Typical first transaction]"),
            ("Repeat Behavior", "[Subscription / Seasonal / Loyal]"),
            ("Referral Potential", "[High / Medium / Low]"),
        ]

    card_w = Inches(4.3)
    card_h = Inches(3.8)
    gap = Inches(0.2)
    y = Inches(1.1)

    build_data_card(slide, MARGIN, y, card_w, card_h, "", left_title, left_rows)
    build_data_card(slide, MARGIN + card_w + gap, y, card_w, card_h, "", right_title, right_rows)

    add_footer(slide)


def slide_07_scenario(prs, variant):
    """Scenario-Based Profile (Before/After)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_header_bar(slide, "SCENARIO-BASED PROFILE")

    card_w = Inches(4.3)
    card_h = Inches(3.5)
    gap = Inches(0.2)
    y = Inches(1.1)

    before_label = "BEFORE Your Solution" if variant == "B2B" else "BEFORE Your Product"
    after_label = "AFTER Your Solution" if variant == "B2B" else "AFTER Your Product"

    # Before card
    left_card = add_card(slide, MARGIN, y, card_w, card_h)
    # Header bar for card
    before_bar = add_rect(slide, MARGIN + Pt(1), y + Pt(1), card_w - Pt(2), Pt(32), fill_color=GREEN)
    add_text(slide, MARGIN + Pt(12), y + Pt(5), card_w - Pt(24), Pt(24),
             before_label, font_size=11, bold=True, color=WHITE)

    if variant == "B2B":
        before_rows = [
            ("Situation", "[Where are they when the problem hits?]"),
            ("Desired Outcome", "[What do they actually want?]"),
            ("Current Approach", "[How are they trying to solve it?]"),
            ("What Goes Wrong", "[Why does it frustrate them?]"),
        ]
    else:
        before_rows = [
            ("Situation", "[What's going on in their life?]"),
            ("Desired Outcome", "[What do they wish they had?]"),
            ("Current Approach", "[How are they handling it now?]"),
            ("Frustration", "[Why does it feel unsolved?]"),
        ]

    row_y = y + Pt(42)
    inner_left = MARGIN + Pt(12)
    inner_w = card_w - Pt(24)
    label_w = Inches(1.4)
    value_w = inner_w - label_w
    row_h = Pt(46)

    for i, (label, value) in enumerate(before_rows):
        add_label_value_row(slide, inner_left, row_y, label_w, value_w, row_h, label, value, is_alt=(i % 2 == 1))
        row_y += row_h

    # After card
    right_card = add_card(slide, MARGIN + card_w + gap, y, card_w, card_h)
    after_bar = add_rect(slide, MARGIN + card_w + gap + Pt(1), y + Pt(1), card_w - Pt(2), Pt(32), fill_color=BRIGHT)
    add_text(slide, MARGIN + card_w + gap + Pt(12), y + Pt(5), card_w - Pt(24), Pt(24),
             after_label, font_size=11, bold=True, color=DARK)

    if variant == "B2B":
        after_rows = [
            ("New Experience", "[How does their business change?]"),
            ("What Makes It Work", "[Capabilities that enable improvement]"),
            ("The Payoff", "[Quantifiable business improvement]"),
        ]
    else:
        after_rows = [
            ("New Experience", "[How their life improves]"),
            ("Emotional Shift", "[How they feel differently]"),
            ("The Result", "[Tangible outcome they can point to]"),
        ]

    row_y = y + Pt(42)
    inner_left_r = MARGIN + card_w + gap + Pt(12)
    row_h_r = Pt(58)

    for i, (label, value) in enumerate(after_rows):
        add_label_value_row(slide, inner_left_r, row_y, label_w, value_w, row_h_r, label, value, is_alt=(i % 2 == 1))
        row_y += row_h_r

    add_footer(slide)


def slide_08_qualifying(prs, variant):
    """Qualifying & Disqualifying Criteria/Signals."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    if variant == "B2B":
        add_header_bar(slide, "QUALIFYING & DISQUALIFYING CRITERIA")
        left_title = "Qualifiers"
        left_rows = [
            ("Must-Haves", "[Non-negotiable criteria]"),
            ("Nice-to-Haves", "[Preferred but not required]"),
            ("Minimum Budget", "[Below this, engagement doesn't make sense]"),
        ]
        right_title = "Disqualifiers"
        right_rows = [
            ("Disqualifiers", "[What makes a company NOT a fit?]"),
            ("Red Flags", "[Warning signs of a bad-fit client]"),
            ("Deal Killers", "[Conditions that end the conversation]"),
        ]
    else:
        add_header_bar(slide, "QUALIFYING & DISQUALIFYING SIGNALS")
        left_title = "Ideal Customer Signals"
        left_rows = [
            ("Behavioral Signals", "[Actions that indicate high intent]"),
            ("Demographic Fit", "[Income, location, life stage match]"),
            ("Engagement Level", "[Email opens, social follows, repeat visits]"),
        ]
        right_title = "Not-a-Fit Signals"
        right_rows = [
            ("Misalignment", "[Needs don't match your offering]"),
            ("Red Flags", "[High return rate, price-only shoppers]"),
            ("Low LTV Indicators", "[One-time buyers, coupon chasers]"),
        ]

    card_w = Inches(4.3)
    card_h = Inches(3.0)
    gap = Inches(0.2)
    y = Inches(1.1)

    build_data_card(slide, MARGIN, y, card_w, card_h, "", left_title, left_rows, accent_color=GREEN)
    build_data_card(slide, MARGIN + card_w + gap, y, card_w, card_h, "", right_title, right_rows, accent_color=RGBColor(0xCC, 0x44, 0x44))

    add_footer(slide)


def slide_09_tiering(prs, variant):
    """Customer Tiering Summary with table + summary boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    title = "CUSTOMER TIERING SUMMARY" if variant == "B2B" else "CUSTOMER SEGMENT TIERING"
    add_header_bar(slide, title)

    # Summary table
    table_left = MARGIN
    table_top = Inches(1.1)
    table_w = CONTENT_W
    cols = 4
    rows = 4

    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_w, Inches(1.8))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(1.7)

    headers = ["Tier", "Segment", "Strategic Approach", "Priority"]
    tier_data = [
        ["Tier 1", "[Type A / Segment A]", "[Maximum investment, direct outreach]", "HIGH"],
        ["Tier 2", "[Type B / Segment B]", "[Scaled outreach, good fit]", "MEDIUM"],
        ["Tier 3", "[Type C / Segment C]", "[Monitor and nurture]", "LOW"],
    ]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN

    for r, row_data in enumerate(tier_data):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = DARK
                paragraph.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = ALTROW if r % 2 == 0 else WHITE

    # Tier 2 & 3 summary cards
    card_w = Inches(4.3)
    card_h = Inches(1.5)
    gap = Inches(0.2)
    y = Inches(3.5)

    t2_desc = "[Customer Type B \u2014 brief description. Why Tier 2? Future potential?]" if variant == "B2B" else "[Segment B \u2014 brief description. Why Tier 2? Nurture potential?]"
    t3_desc = "[Customer Type C \u2014 brief description. Why Tier 3? Worth monitoring?]" if variant == "B2B" else "[Segment C \u2014 brief description. Why Tier 3? Worth monitoring?]"

    # Tier 2 card
    t2_card = add_card(slide, MARGIN, y, card_w, card_h, accent_color=BRIGHT)
    add_text(slide, MARGIN + Pt(15), y + Pt(8), card_w - Pt(30), Pt(20),
             "Tier 2 Summary", font_size=10, bold=True, color=GREEN)
    add_text(slide, MARGIN + Pt(15), y + Pt(32), card_w - Pt(30), Inches(0.8),
             t2_desc, font_size=9, color=DARK)

    # Tier 3 card
    t3_card = add_card(slide, MARGIN + card_w + gap, y, card_w, card_h, accent_color=BORDER)
    add_text(slide, MARGIN + card_w + gap + Pt(15), y + Pt(8), card_w - Pt(30), Pt(20),
             "Tier 3 Summary", font_size=10, bold=True, color=GREEN)
    add_text(slide, MARGIN + card_w + gap + Pt(15), y + Pt(32), card_w - Pt(30), Inches(0.8),
             t3_desc, font_size=9, color=DARK)

    add_footer(slide)


def slide_10_persona(prs, variant):
    """Primary Buyer Persona."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    add_header_bar(slide, "PRIMARY BUYER PERSONA")

    # Profile header card
    profile_card = add_card(slide, MARGIN, Inches(1.05), CONTENT_W, Inches(0.65))
    name_text = '[Persona Name \u2014 e.g., "The Business Owner"]' if variant == "B2B" else '[Persona Name \u2014 e.g., "Busy Mom Beth"]'
    subtitle_text = "[Champion / Decision Maker]  |  [Title]  |  Reports to [Title]" if variant == "B2B" else "[Age]  |  [Location]  |  [Occupation]  |  [Household]"
    add_text(slide, MARGIN + Pt(15), Inches(1.12), CONTENT_W - Pt(30), Pt(22),
             name_text, font_size=13, bold=True, color=DARK)
    add_text(slide, MARGIN + Pt(15), Inches(1.42), CONTENT_W - Pt(30), Pt(18),
             subtitle_text, font_size=9, color=DARK)

    # 4 quadrant cards
    card_w = Inches(4.3)
    card_h = Inches(1.35)
    gap = Inches(0.2)
    row1_y = Inches(1.9)
    row2_y = Inches(3.4)

    if variant == "B2B":
        q1_title = "Goals & Motivations"
        q1_rows = [("Primary Goal", "[What they're trying to achieve]"), ("Success Metric", "[How they measure a win]")]
        q2_title = "Pain Points"
        q2_rows = [("Frustration", "[Keeps them up at night]"), ("What's Blocking", "[Why they haven't solved it]")]
        q3_title = "How They Decide"
        q3_rows = [("Research", "[Google, LinkedIn, peers, etc.]"), ("Trust Signals", "[Case studies, ROI data, referrals]")]
        q4_title = "How to Reach Them"
        q4_rows = [("Best Channels", "[Where marketing reaches them]"), ("Content Type", "[Video, case studies, whitepapers]")]
    else:
        q1_title = "Goals & Motivations"
        q1_rows = [("Primary Goal", "[What they want to achieve or feel]"), ("Success Looks Like", "[Their ideal outcome]")]
        q2_title = "Frustrations"
        q2_rows = [("Main Pain", "[What keeps them stuck]"), ("Emotional Toll", "[How the problem feels]")]
        q3_title = "How They Discover"
        q3_rows = [("Channels", "[Instagram, Google, friends, TikTok]"), ("Content Preferences", "[Video, reviews, tutorials, UGC]")]
        q4_title = "How to Convert Them"
        q4_rows = [("Best Offer", "[Free trial, discount, bundle]"), ("Trust Builder", "[Reviews, guarantees, social proof]")]

    for idx, (title, rows, x, y) in enumerate([
        (q1_title, q1_rows, MARGIN, row1_y),
        (q2_title, q2_rows, MARGIN + card_w + gap, row1_y),
        (q3_title, q3_rows, MARGIN, row2_y),
        (q4_title, q4_rows, MARGIN + card_w + gap, row2_y),
    ]):
        card = add_card(slide, x, y, card_w, card_h)
        add_text(slide, x + Pt(15), y + Pt(6), card_w - Pt(30), Pt(18),
                 title, font_size=10, bold=True, color=GREEN)
        add_divider(slide, x + Pt(15), y + Pt(26), card_w - Pt(30), thickness=Pt(1))

        row_y = y + Pt(32)
        label_w = Inches(1.4)
        value_w = card_w - Pt(30) - label_w
        for i, (label, value) in enumerate(rows):
            add_label_value_row(slide, x + Pt(15), row_y, label_w, value_w, Pt(28), label, value, is_alt=(i % 2 == 1))
            row_y += Pt(28)

    add_footer(slide)


def slide_11_next_steps(prs, variant):
    """Next Steps slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)

    # Top accent
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=GREEN)

    add_text(slide, MARGIN, Inches(0.8), CONTENT_W, Inches(0.6),
             "NEXT STEPS", font_size=28, bold=True, color=DARK)

    add_divider(slide, MARGIN, Inches(1.5), Inches(2.5), color=GREEN, thickness=Pt(2))

    steps_card = add_card(slide, MARGIN, Inches(1.8), CONTENT_W, Inches(2.5))

    steps = [
        "1.  Review and validate ICP findings with your team",
        "2.  Proceed to Stage 2: Market Research & Competitive Analysis",
        "3.  Use Tier 1 insights to guide keyword and competitor research",
        "4.  Identify competitive gaps and positioning opportunities",
        "5.  Prepare for Stage 3: Value Proposition & Messaging Development",
    ]

    y = Inches(1.9)
    for i, step in enumerate(steps):
        bg_color = ALTROW if i % 2 == 0 else None
        if bg_color:
            add_rect(slide, MARGIN + Pt(10), y, CONTENT_W - Pt(20), Pt(28), fill_color=bg_color)
        add_text(slide, MARGIN + Pt(20), y + Pt(4), CONTENT_W - Pt(40), Pt(20),
                 step, font_size=11, color=DARK)
        y += Pt(30)

    add_text(slide, MARGIN, Inches(4.5), CONTENT_W, Inches(0.4),
             "Questions? Let's discuss.", font_size=14, bold=True, color=GREEN)

    # Footer
    add_rect(slide, Inches(0), Inches(5.33), SLIDE_W, Inches(0.295), fill_color=DARK)
    add_text(slide, MARGIN, Inches(5.35), Inches(9), Inches(0.25),
             "ClearThink Marketing", font_size=8, color=LGGREEN)


# ═══════════════════════════════════════════════════
#  BUILD DECKS
# ═══════════════════════════════════════════════════

def build_deck(variant, output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs, variant)
    slide_02_tiers(prs, variant)
    slide_03_deep_dive(prs, variant)
    slide_04_profile_details(prs, variant)
    slide_05_pain_triggers(prs, variant)
    slide_06_landscape(prs, variant)
    slide_07_scenario(prs, variant)
    slide_08_qualifying(prs, variant)
    slide_09_tiering(prs, variant)
    slide_10_persona(prs, variant)
    slide_11_next_steps(prs, variant)

    prs.save(output_path)
    print(f"Built: {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    import os
    base = os.path.dirname(os.path.abspath(__file__))

    build_deck("B2B", os.path.join(base, "ClearLaunch_B2B_ICP_Summary_Deck.pptx"))
    build_deck("B2C", os.path.join(base, "ClearLaunch_B2C_ICP_Summary_Deck.pptx"))
    print("\nDone. Both decks rebuilt.")
