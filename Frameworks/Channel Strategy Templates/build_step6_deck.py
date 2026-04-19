"""
Build ClearLaunch Step 6 (Channel Strategy & Customer Journey) Summary Deck (.pptx)
8 slides with ClearThink brand styling. Industry-agnostic version.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
GREEN = RGBColor(0x1B, 0x9B, 0x5E)
BRIGHT = RGBColor(0x3B, 0xEB, 0x96)
DARK = RGBColor(0x12, 0x17, 0x18)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF6, 0xF3, 0xEF)
LGGREEN = RGBColor(0xE5, 0xF5, 0xEC)
ALTROW = RGBColor(0xF0, 0xF9, 0xF4)
BORDER = RGBColor(0xA8, 0xD9, 0xBD)
ORANGE = RGBColor(0xD4, 0x88, 0x0F)
LIGHTGRAY = RGBColor(0xF0, 0xF0, 0xF0)
MIDGRAY = RGBColor(0x66, 0x66, 0x66)
DARKGREEN = RGBColor(0x15, 0x7A, 0x49)
DARKGREEN_BG = RGBColor(0x1E, 0x2A, 0x20)

# Slide dimensions
SLIDE_W = 9144000
SLIDE_H = 5143500

# Common measurements
MARGIN_L = Emu(548640)
CONTENT_W = Emu(8046720)
HEADER_H = Emu(777240)
FOOTER_Y = Emu(4873752)
FOOTER_H = Emu(269748)


# ========================================
# REUSABLE HELPERS (from build_offer_dev_deck)
# ========================================

def add_header_banner(slide, title_text):
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    banner.fill.solid()
    banner.fill.fore_color.rgb = GREEN
    banner.line.fill.background()
    title = slide.shapes.add_textbox(MARGIN_L, Emu(137160), CONTENT_W, Emu(502920))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Calibri"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide, page_num=None, total=8, client="[Client Name]"):
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, FOOTER_Y, SLIDE_W, FOOTER_H)
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = DARK
    footer_bar.line.fill.background()
    # Left text
    left_box = slide.shapes.add_textbox(MARGIN_L, Emu(4892040), Emu(3000000), Emu(228600))
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CLEARTHINK MARKETING"
    p.font.name = "Calibri"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = BRIGHT
    # Right text
    if page_num:
        right_box = slide.shapes.add_textbox(Emu(5500000), Emu(4892040), Emu(3200000), Emu(228600))
        tf2 = right_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        p2.text = f"ClearLaunch Step 6  |  {client}  |  {page_num}/{total}"
        p2.font.name = "Calibri"
        p2.font.size = Pt(8)
        p2.font.color.rgb = MIDGRAY


def add_content_card(slide, left, top, width, height, label, placeholder,
                     accent_color=GREEN):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)
    # Left accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Emu(top + 25400),
                                     Emu(50800), Emu(height - 50800))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    # Label
    label_box = slide.shapes.add_textbox(Emu(left + 190500), Emu(top + 101600),
                                          Emu(width - 254000), Emu(228600))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = accent_color
    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left + 190500),
                                      Emu(top + 330200), Emu(width - 381000), Emu(19050))
    divider.fill.solid()
    divider.fill.fore_color.rgb = LGGREEN
    divider.line.fill.background()
    # Placeholder text
    ph_box = slide.shapes.add_textbox(Emu(left + 190500), Emu(top + 406400),
                                       Emu(width - 381000), Emu(height - 508000))
    tf2 = ph_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = placeholder
    p2.font.name = "Calibri"
    p2.font.size = Pt(10)
    p2.font.color.rgb = DARK


# ========================================
# NEW HELPERS FOR STEP 6
# ========================================

def add_badge_pill(slide, left, top, text, color):
    """Small pill-shaped badge (e.g., DATA-BACKED ESTIMATE)."""
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   Emu(2194560), Emu(254000))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_stat_card(slide, left, top, width, height, value, label, sublabel,
                  dark_bg=False, value_color=GREEN):
    """Big-number stat card for investment overview."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = DARK if dark_bg else WHITE
    card.line.color.rgb = BORDER if not dark_bg else DARK
    card.line.width = Pt(0.5)
    # Value
    val_box = slide.shapes.add_textbox(left, Emu(top + 152400), width, Emu(406400))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = value
    p.font.name = "Calibri"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = value_color
    # Label
    lbl_box = slide.shapes.add_textbox(left, Emu(top + 558800), width, Emu(228600))
    tf2 = lbl_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    p2.font.name = "Calibri"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = WHITE if dark_bg else DARK
    # Sublabel
    sub_box = slide.shapes.add_textbox(left, Emu(top + 787400), width, Emu(228600))
    tf3 = sub_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.text = sublabel
    p3.font.name = "Calibri"
    p3.font.size = Pt(10)
    p3.font.color.rgb = MIDGRAY


def add_table_shape(slide, left, top, headers, rows, col_widths, row_height=304800):
    """Create a table as shapes (rectangles with text) for precise styling."""
    total_w = sum(col_widths)
    header_h = Emu(304800)

    # Header row
    x = left
    for j, h in enumerate(headers):
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top,
                                       Emu(col_widths[j]), header_h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        cell.line.fill.background()
        tf = cell.text_frame
        tf.margin_left = Emu(76200)
        tf.margin_top = Emu(50800)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = h
        p.font.name = "Calibri"
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = WHITE
        x += Emu(col_widths[j])

    # Data rows
    for i, row_data in enumerate(rows):
        y = Emu(top + int(header_h) + i * row_height)
        x = left
        fill = ALTROW if i % 2 == 0 else WHITE
        for j, text in enumerate(row_data):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                           Emu(col_widths[j]), Emu(row_height))
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            cell.line.fill.background()
            tf = cell.text_frame
            tf.margin_left = Emu(76200)
            tf.margin_top = Emu(50800)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.name = "Calibri"
            p.font.size = Pt(8)
            p.font.color.rgb = DARK
            x += Emu(col_widths[j])


def add_stats_panel(slide, left, top, stats, accent_color=GREEN):
    """Right-side stats panel for projection slides."""
    panel_w = Emu(2926080)
    panel_h = Emu(3291840)
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, panel_w, panel_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHTGRAY
    bg.line.fill.background()
    # Title
    title_box = slide.shapes.add_textbox(Emu(left + 152400), Emu(top + 101600),
                                          Emu(int(panel_w) - 304800), Emu(228600))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "COMBINED PERFORMANCE"
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = accent_color
    # Stats
    y_offset = top + 406400
    for value, label in stats:
        # Value
        val_box = slide.shapes.add_textbox(Emu(left + 152400), Emu(y_offset),
                                            Emu(int(panel_w) - 304800), Emu(304800))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK
        # Label
        lbl_box = slide.shapes.add_textbox(Emu(left + 152400), Emu(y_offset + 254000),
                                            Emu(int(panel_w) - 304800), Emu(228600))
        tf2 = lbl_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = label
        p2.font.name = "Calibri"
        p2.font.size = Pt(9)
        p2.font.color.rgb = MIDGRAY
        y_offset += 533400


# ========================================
# BUILD THE DECK
# ========================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ========================================
# SLIDE 1: Title
# ========================================
slide = prs.slides.add_slide(blank)

# Thin top bar
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(50800))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = GREEN
top_bar.line.fill.background()

# Title
title = slide.shapes.add_textbox(MARGIN_L, Emu(1005840), CONTENT_W, Emu(457200))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "CHANNEL STRATEGY"
p.font.name = "Calibri"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK

# Subtitle in green
sub1 = slide.shapes.add_textbox(MARGIN_L, Emu(1463040), CONTENT_W, Emu(457200))
tf = sub1.text_frame
p = tf.paragraphs[0]
p.text = "& CUSTOMER JOURNEY"
p.font.name = "Calibri"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = GREEN

# Client name
client = slide.shapes.add_textbox(MARGIN_L, Emu(2133600), CONTENT_W, Emu(381000))
tf = client.text_frame
p = tf.paragraphs[0]
p.text = "[Client Name]"
p.font.name = "Calibri"
p.font.size = Pt(20)
p.font.color.rgb = DARK

# Divider
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(2667000),
                                  Emu(2286000), Emu(25400))
divider.fill.solid()
divider.fill.fore_color.rgb = BRIGHT
divider.line.fill.background()

# Tagline
tag = slide.shapes.add_textbox(MARGIN_L, Emu(2819400), CONTENT_W, Emu(304800))
tf = tag.text_frame
p = tf.paragraphs[0]
p.text = "[Industry Context]  |  [Location]"
p.font.name = "Calibri"
p.font.size = Pt(13)
p.font.color.rgb = MIDGRAY

# Step label + date
step_lbl = slide.shapes.add_textbox(MARGIN_L, Emu(3200400), CONTENT_W, Emu(381000))
tf = step_lbl.text_frame
p = tf.paragraphs[0]
p.text = "ClearLaunch System  |  Step 6"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GREEN

date_box = slide.shapes.add_textbox(MARGIN_L, Emu(3505200), CONTENT_W, Emu(228600))
tf = date_box.text_frame
p = tf.paragraphs[0]
p.text = "[Date]"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = MIDGRAY

add_footer(slide, page_num=1)

# ========================================
# SLIDE 2: The Customer Journey
# ========================================
slide = prs.slides.add_slide(blank)
add_header_banner(slide, "THE CUSTOMER JOURNEY")

# Subtitle
sub = slide.shapes.add_textbox(MARGIN_L, Emu(850000), CONTENT_W, Emu(228600))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "[Number] stages from first awareness to [outcome]"
p.font.name = "Calibri"
p.font.size = Pt(13)
p.font.color.rgb = MIDGRAY

# 4 journey stage cards
card_w = Emu(1962150)
card_h = Emu(2926080)
card_gap = Emu(101600)
card_top = Emu(1143000)

stage_colors = [BRIGHT, GREEN, DARKGREEN, DARK]
stage_names = ["TOFU\n(Awareness)", "MOFU\n(Consideration)", "BOFU\n(Decision)", "PURCHASE"]

for i, (color, name) in enumerate(zip(stage_colors, stage_names)):
    x = int(MARGIN_L) + i * (int(card_w) + int(card_gap))
    y = int(card_top)

    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)

    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Emu(57150))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    # Numbered circle
    circle_d = Emu(508000)
    cx = x + (int(card_w) - int(circle_d)) // 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, Emu(y + 152400),
                                     circle_d, circle_d)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(i + 1)
    p.font.name = "Calibri"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Stage name
    name_box = slide.shapes.add_textbox(x, Emu(y + 736600), card_w, Emu(406400))
    tf = name_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = name
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = color

    # Description
    desc_box = slide.shapes.add_textbox(Emu(x + 76200), Emu(y + 1219200),
                                         Emu(int(card_w) - 152400), Emu(457200))
    tf2 = desc_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "[Customer mindset at this stage]"
    p2.font.name = "Calibri"
    p2.font.size = Pt(9)
    p2.font.color.rgb = DARK

    # Channels label
    ch_lbl = slide.shapes.add_textbox(Emu(x + 76200), Emu(y + 1828800),
                                       Emu(int(card_w) - 152400), Emu(152400))
    tf3 = ch_lbl.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.text = "CHANNELS"
    p3.font.name = "Calibri"
    p3.font.size = Pt(8)
    p3.font.color.rgb = MIDGRAY

    # Channels value
    ch_val = slide.shapes.add_textbox(Emu(x + 76200), Emu(y + 1981200),
                                       Emu(int(card_w) - 152400), Emu(304800))
    tf4 = ch_val.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    p4.text = "[Active channels]"
    p4.font.name = "Calibri"
    p4.font.size = Pt(9)
    p4.font.bold = True
    p4.font.color.rgb = GREEN

    # KPI label
    kpi_lbl = slide.shapes.add_textbox(Emu(x + 76200), Emu(y + 2438400),
                                        Emu(int(card_w) - 152400), Emu(152400))
    tf5 = kpi_lbl.text_frame
    p5 = tf5.paragraphs[0]
    p5.alignment = PP_ALIGN.CENTER
    p5.text = "KEY KPI"
    p5.font.name = "Calibri"
    p5.font.size = Pt(8)
    p5.font.color.rgb = MIDGRAY

    # KPI value
    kpi_val = slide.shapes.add_textbox(Emu(x + 76200), Emu(y + 2590800),
                                        Emu(int(card_w) - 152400), Emu(228600))
    tf6 = kpi_val.text_frame
    p6 = tf6.paragraphs[0]
    p6.alignment = PP_ALIGN.CENTER
    p6.text = "[Primary metric]"
    p6.font.name = "Calibri"
    p6.font.size = Pt(9)
    p6.font.color.rgb = DARK

    # Arrow between cards (except after last)
    if i < 3:
        arrow_x = x + int(card_w) + int(card_gap) // 2 - Emu(76200)
        arrow_y = y + int(card_h) // 2
        arrow = slide.shapes.add_textbox(arrow_x, arrow_y, Emu(152400), Emu(228600))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = "\u25B8"
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = BRIGHT

add_footer(slide, page_num=2)

# ========================================
# SLIDE 3: Channel Universe
# ========================================
slide = prs.slides.add_slide(blank)
add_header_banner(slide, "CHANNEL UNIVERSE")

sub = slide.shapes.add_textbox(MARGIN_L, Emu(850000), CONTENT_W, Emu(228600))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "7 channels evaluated against [Client Name]'s ICP, budget, and data availability"
p.font.name = "Calibri"
p.font.size = Pt(13)
p.font.color.rgb = MIDGRAY

# Channel matrix table
col_widths = [1371600, 1143000, 1016000, 914400, 762000, 1143000]
add_table_shape(slide, int(MARGIN_L), 1143000,
    headers=["Channel", "Traffic Type", "Stage", "Data Avail.", "ICP Fit", "Recommendation"],
    rows=[
        ["Google Ads", "Paid / Intent", "MOFU-BOFU", "High", "[score]", "[tier]"],
        ["Meta Ads", "Paid / Interrupt", "TOFU-MOFU", "Moderate", "[score]", "[tier]"],
        ["LinkedIn Ads", "Paid / Prof.", "TOFU-MOFU", "Moderate", "[score]", "[tier]"],
        ["TikTok", "Paid / Interrupt", "TOFU", "Low", "[score]", "[tier]"],
        ["SEO", "Organic / Intent", "TOFU-BOFU", "High", "[score]", "[tier]"],
        ["Content Mktg", "Organic / Edu", "TOFU-MOFU", "Moderate", "[score]", "[tier]"],
        ["Organic Social", "Organic / Comm", "TOFU", "Low", "[score]", "[tier]"],
    ],
    col_widths=col_widths, row_height=254000)

# Decision rules callout
rules_y = Emu(3657600)
rules_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, rules_y,
                                    CONTENT_W, Emu(1066800))
rules_bg.fill.solid()
rules_bg.fill.fore_color.rgb = LIGHTGRAY
rules_bg.line.fill.background()

rules_lbl = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 152400), Emu(int(rules_y) + 76200),
                                      Emu(int(CONTENT_W) - 304800), Emu(228600))
tf = rules_lbl.text_frame
p = tf.paragraphs[0]
p.text = "DECISION RULES"
p.font.name = "Calibri"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = GREEN

rules_txt = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 152400), Emu(int(rules_y) + 304800),
                                      Emu(int(CONTENT_W) - 304800), Emu(685800))
tf2 = rules_txt.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = ("1. ICP Presence Gate \u2014 [rationale]\n"
           "2. Budget Viability Gate \u2014 [rationale]\n"
           "3. Journey Coverage Check \u2014 [rationale]")
p2.font.name = "Calibri"
p2.font.size = Pt(9)
p2.font.color.rgb = DARK

add_footer(slide, page_num=3)

# ========================================
# SLIDE 4: Channel Recommendation
# ========================================
slide = prs.slides.add_slide(blank)
add_header_banner(slide, "CHANNEL RECOMMENDATION")

sub = slide.shapes.add_textbox(MARGIN_L, Emu(850000), CONTENT_W, Emu(228600))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "[N] primary channels + [N] supporting channel"
p.font.name = "Calibri"
p.font.size = Pt(13)
p.font.color.rgb = MIDGRAY

# 3 horizontal channel cards
card_configs = [
    (GREEN, "PRIMARY", "[Primary Channel 1]", "[Role Title]",
     "[2-3 sentences explaining channel role and funnel stages served]"),
    (GREEN, "PRIMARY", "[Primary Channel 2]", "[Role Title]",
     "[2-3 sentences explaining channel role and funnel stages served]"),
    (ORANGE, "SUPPORTING", "[Supporting Channel]", "[Role Title]",
     "[2-3 sentences explaining channel role and why supporting vs primary]"),
]

card_top = 1143000
card_h = Emu(914400)
card_gap = Emu(101600)

for i, (color, tier, name, role, desc) in enumerate(card_configs):
    y = card_top + i * (int(card_h) + int(card_gap))

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, y,
                                   CONTENT_W, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)

    # Left accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(y + 25400),
                                     Emu(50800), Emu(int(card_h) - 50800))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    # Tier badge
    badge_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), Emu(y + 76200),
                                          Emu(1524000), Emu(228600))
    tf = badge_box.text_frame
    p = tf.paragraphs[0]
    p.text = tier
    p.font.name = "Calibri"
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = color

    # Channel name + role
    name_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), Emu(y + 254000),
                                         Emu(int(CONTENT_W) - 381000), Emu(304800))
    tf = name_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = name
    r1.font.name = "Calibri"
    r1.font.size = Pt(16)
    r1.font.bold = True
    r1.font.color.rgb = DARK
    r2 = p.add_run()
    r2.text = f"  \u2014  {role}"
    r2.font.name = "Calibri"
    r2.font.size = Pt(16)
    r2.font.italic = True
    r2.font.color.rgb = color

    # Description
    desc_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), Emu(y + 558800),
                                         Emu(int(CONTENT_W) - 381000), Emu(int(card_h) - 635000))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = MIDGRAY

add_footer(slide, page_num=4)

# ========================================
# SLIDE 5: Investment Overview
# ========================================
slide = prs.slides.add_slide(blank)
add_header_banner(slide, "INVESTMENT OVERVIEW")

sub = slide.shapes.add_textbox(MARGIN_L, Emu(850000), CONTENT_W, Emu(228600))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "Retainers are fixed. Ad spend is the control variable."
p.font.name = "Calibri"
p.font.size = Pt(13)
p.font.color.rgb = MIDGRAY

# 3 stat cards
stat_w = Emu(2614930)
stat_h = Emu(1143000)
stat_gap = Emu(101600)
stat_top = Emu(1143000)

add_stat_card(slide, int(MARGIN_L), int(stat_top), stat_w, stat_h,
              "[retainer total]", "Monthly Retainers", "Fixed across channels")
add_stat_card(slide, int(MARGIN_L) + int(stat_w) + int(stat_gap), int(stat_top),
              stat_w, stat_h,
              "[ad spend total]", "Monthly Ad Spend", "Adjustable control variable")
add_stat_card(slide, int(MARGIN_L) + 2 * (int(stat_w) + int(stat_gap)), int(stat_top),
              stat_w, stat_h,
              "[grand total]", "Total Monthly", "All-in investment",
              dark_bg=True, value_color=BRIGHT)

# Investment table below
table_top = int(stat_top) + int(stat_h) + 228600
col_widths = [2133600, 1524000, 1524000, 1524000]
add_table_shape(slide, int(MARGIN_L), table_top,
    headers=["Channel", "Retainer", "Ad Spend", "Total"],
    rows=[
        ["[Channel 1]", "[amount]", "[amount]", "[amount]"],
        ["[Channel 2]", "[amount]", "[amount]", "[amount]"],
        ["[Channel 3]", "[amount]", "[amount]", "[amount]"],
    ],
    col_widths=col_widths, row_height=254000)

add_footer(slide, page_num=5)

# ========================================
# SLIDE 6: Projections — Data-Backed
# ========================================
slide = prs.slides.add_slide(blank)
add_header_banner(slide, "PROJECTIONS: [DATA-BACKED CHANNEL]")

add_badge_pill(slide, int(MARGIN_L), 850000, "DATA-BACKED ESTIMATE", GREEN)

sub = slide.shapes.add_textbox(MARGIN_L, Emu(1143000), Emu(5200000), Emu(228600))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "Built from real keyword data (Step 3 Market Research)"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = MIDGRAY

# Left: projection table
proj_col_widths = [1524000, 1143000, 1143000]
add_table_shape(slide, int(MARGIN_L), 1371600,
    headers=["Metric", "[Line 1]", "[Line 2]"],
    rows=[
        ["Monthly Budget", "[amount]", "[amount]"],
        ["Avg CPC", "[amount]", "[amount]"],
        ["Monthly Clicks", "[calc]", "[calc]"],
        ["Conv Rate", "[rate]", "[rate]"],
        ["Monthly Leads", "[calc]", "[calc]"],
        ["Close Rate", "[rate]", "[rate]"],
        ["New Clients", "[calc]", "[calc]"],
        ["Revenue", "[calc]", "[calc]"],
        ["Ad Spend ROAS", "[calc]", "[calc]"],
        ["Profit", "[calc]", "[calc]"],
    ],
    col_widths=proj_col_widths, row_height=254000)

# Right: stats panel
add_stats_panel(slide, int(MARGIN_L) + 4114800, 1371600,
    stats=[
        ("[total revenue]", "Monthly Revenue"),
        ("[total clients]", "New Clients/Month"),
        ("[true ROAS]", "True ROAS (incl. retainer)"),
        ("[CPA]", "Cost Per Acquisition"),
        ("[margin]", "Profit Margin"),
    ],
    accent_color=GREEN)

add_footer(slide, page_num=6)

# ========================================
# SLIDE 7: Projections — Modeled
# ========================================
slide = prs.slides.add_slide(blank)
add_header_banner(slide, "PROJECTIONS: [MODELED CHANNEL]")

add_badge_pill(slide, int(MARGIN_L), 850000, "MODELED ESTIMATE", ORANGE)

sub = slide.shapes.add_textbox(MARGIN_L, Emu(1143000), Emu(5200000), Emu(228600))
tf = sub.text_frame
p = tf.paragraphs[0]
p.text = "Built from industry CPM/CTR benchmarks"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = MIDGRAY

# Left: projection table
add_table_shape(slide, int(MARGIN_L), 1371600,
    headers=["Metric", "[Line 1]", "[Line 2]"],
    rows=[
        ["Monthly Budget", "[amount]", "[amount]"],
        ["CPM", "[amount]", "[amount]"],
        ["CTR", "[rate]", "[rate]"],
        ["Impressions", "[calc]", "[calc]"],
        ["Clicks", "[calc]", "[calc]"],
        ["Conv Rate", "[rate]", "[rate]"],
        ["Leads", "[calc]", "[calc]"],
        ["New Clients", "[calc]", "[calc]"],
        ["Revenue", "[calc]", "[calc]"],
        ["Ad Spend ROAS", "[calc]", "[calc]"],
    ],
    col_widths=proj_col_widths, row_height=254000)

# Right: stats panel (orange accent)
add_stats_panel(slide, int(MARGIN_L) + 4114800, 1371600,
    stats=[
        ("[total revenue]", "Monthly Revenue"),
        ("[total clients]", "New Clients/Month"),
        ("[true ROAS]", "True ROAS (incl. retainer)"),
        ("[CPA]", "Cost Per Acquisition"),
        ("[margin]", "Profit Margin"),
    ],
    accent_color=ORANGE)

# Benchmark citation
cite = slide.shapes.add_textbox(MARGIN_L, Emu(4648200), CONTENT_W, Emu(152400))
tf = cite.text_frame
p = tf.paragraphs[0]
p.text = "Benchmarks: [Source 1], [Source 2], [Source 3]"
p.font.name = "Calibri"
p.font.size = Pt(8)
p.font.italic = True
p.font.color.rgb = MIDGRAY

add_footer(slide, page_num=7)

# ========================================
# SLIDE 8: Key Takeaways & Next Steps
# ========================================
slide = prs.slides.add_slide(blank)

# Dark background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

# Left accent bar
accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(109728), SLIDE_H)
accent.fill.solid()
accent.fill.fore_color.rgb = GREEN
accent.line.fill.background()

# Title
title = slide.shapes.add_textbox(MARGIN_L, Emu(381000), CONTENT_W, Emu(457200))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "KEY TAKEAWAYS & NEXT STEPS"
p.font.name = "Calibri"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Divider
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(914400),
                                  Emu(2286000), Emu(25400))
divider.fill.solid()
divider.fill.fore_color.rgb = GREEN
divider.line.fill.background()

# 5 takeaway bullets
takeaways = [
    "[Primary channel role and why it's the foundation]",
    "[Secondary channel role and how it complements]",
    "[Supporting channel for demand generation]",
    "[How the funnel architecture works across channels]",
    "[Budget control lever and scaling strategy]",
]

bullet_y = 1143000
for i, text in enumerate(takeaways):
    y = bullet_y + i * 533400
    # Green circle bullet
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN_L, Emu(y + 38100),
                                     Emu(228600), Emu(228600))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    # Takeaway text
    txt = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 381000), Emu(y),
                                    Emu(7620000), Emu(381000))
    tf = txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

# "NEXT: STEP 7" callout box
next_y = Emu(3886200)
next_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, next_y,
                                   CONTENT_W, Emu(762000))
next_bg.fill.solid()
next_bg.fill.fore_color.rgb = DARKGREEN_BG
next_bg.line.fill.background()

next_lbl = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 152400), Emu(int(next_y) + 76200),
                                     Emu(int(CONTENT_W) - 304800), Emu(228600))
tf = next_lbl.text_frame
p = tf.paragraphs[0]
p.text = "NEXT: STEP 7 \u2014 SUCCESS METRICS & LAUNCH ROADMAP"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GREEN

next_txt = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 152400), Emu(int(next_y) + 330200),
                                     Emu(int(CONTENT_W) - 304800), Emu(381000))
tf2 = next_txt.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = ('Step 6 answered: "Where should we show up, and what should we expect?"\n'
           'Step 7 will answer: "How do we execute, in what order, and how do we know it\'s working?"')
p2.font.name = "Calibri"
p2.font.size = Pt(10)
p2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)

# Footer (green bar for last slide)
footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, FOOTER_Y, SLIDE_W, FOOTER_H)
footer_bar.fill.solid()
footer_bar.fill.fore_color.rgb = GREEN
footer_bar.line.fill.background()

left_txt = slide.shapes.add_textbox(MARGIN_L, Emu(4892040), Emu(3000000), Emu(228600))
tf = left_txt.text_frame
p = tf.paragraphs[0]
p.text = "CLEARTHINK MARKETING"
p.font.name = "Calibri"
p.font.size = Pt(8)
p.font.bold = True
p.font.color.rgb = WHITE

right_txt = slide.shapes.add_textbox(Emu(5500000), Emu(4892040), Emu(3200000), Emu(228600))
tf2 = right_txt.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.RIGHT
p2.text = "clearthinkmarketing.com"
p2.font.name = "Calibri"
p2.font.size = Pt(8)
p2.font.color.rgb = WHITE


# ========================================
# SAVE
# ========================================
prs.save("Frameworks/ClearLaunch_Step6_ChannelStrategy_Summary_Deck.pptx")
print("Created: Frameworks/ClearLaunch_Step6_ChannelStrategy_Summary_Deck.pptx (8 slides)")
