"""
Build ClearLaunch Offer Dev Summary Deck (.pptx) — 8 slides
Matches the UVP/ICP Summary Deck styling pattern exactly.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Brand colors
GREEN = RGBColor(0x1B, 0x9B, 0x5E)
BRIGHT = RGBColor(0x3B, 0xEB, 0x96)
DARK = RGBColor(0x12, 0x17, 0x18)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF6, 0xF3, 0xEF)
LGGREEN = RGBColor(0xE5, 0xF5, 0xEC)
ALTROW = RGBColor(0xF0, 0xF9, 0xF4)
BORDER = RGBColor(0xA8, 0xD9, 0xBD)

# Slide dimensions (matching UVP/ICP deck)
SLIDE_W = 9144000
SLIDE_H = 5143500

# Common measurements
MARGIN_L = Emu(548640)
CONTENT_W = Emu(8046720)
HEADER_H = Emu(777240)
FOOTER_Y = Emu(4873752)
FOOTER_H = Emu(269748)


def add_header_banner(slide, title_text):
    """Green header banner with white title text."""
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, HEADER_H)
    banner.fill.solid()
    banner.fill.fore_color.rgb = GREEN
    banner.line.fill.background()

    title = slide.shapes.add_textbox(MARGIN_L, Emu(137160), CONTENT_W, Emu(502920))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Calibri"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_footer(slide, subtitle="Offer Dev Summary"):
    """Dark footer bar with ClearThink branding."""
    footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, FOOTER_Y, SLIDE_W, FOOTER_H)
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = DARK
    footer_bar.line.fill.background()

    footer_text = slide.shapes.add_textbox(MARGIN_L, Emu(4892040), Emu(8229600), Emu(228600))
    tf = footer_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"ClearThink Marketing  |  {subtitle}"
    p.font.name = "Calibri"
    p.font.size = Pt(8)
    p.font.color.rgb = LGGREEN


def add_content_card(slide, left, top, width, height, label, placeholder,
                     accent_color=GREEN):
    """Rounded rectangle card with accent bar, label, and placeholder text."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)

    # Left accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Emu(top + 25400), Emu(50800),
                                     Emu(height - 50800))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    # Label
    label_box = slide.shapes.add_textbox(Emu(left + 190500), Emu(top + 101600),
                                          Emu(width - 254000), Emu(228600))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left + 190500),
                                      Emu(top + 330200), Emu(width - 381000), Emu(19050))
    divider.fill.solid()
    divider.fill.fore_color.rgb = LGGREEN
    divider.line.fill.background()

    # Placeholder text
    ph_box = slide.shapes.add_textbox(Emu(left + 190500), Emu(top + 406400),
                                       Emu(width - 381000), Emu(height - 508000))
    tf2 = ph_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = placeholder
    p2.font.name = "Calibri"
    p2.font.size = Pt(10)
    p2.font.color.rgb = DARK


def add_numbered_badge(slide, left, top, number, bg_color=GREEN):
    """Small numbered badge (like tier badges)."""
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(381000), Emu(254000))
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()

    tf = badge.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(number)
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE


# ========================================
# BUILD THE DECK
# ========================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]  # Blank

# ========================================
# SLIDE 1: Title
# ========================================
slide = prs.slides.add_slide(blank_layout)

# Thin top bar
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(50800))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = GREEN
top_bar.line.fill.background()

# Title
title = slide.shapes.add_textbox(MARGIN_L, Emu(1005840), CONTENT_W, Emu(731520))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "OFFER DEVELOPMENT"
p.font.name = "Calibri"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = DARK

# OFFER badge
badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(1828800),
                                Emu(731520), Emu(274320))
badge.fill.solid()
badge.fill.fore_color.rgb = GREEN
badge.line.fill.background()
badge_text = slide.shapes.add_textbox(MARGIN_L, Emu(1828800), Emu(731520), Emu(274320))
tf = badge_text.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "OFFER"
p.font.name = "Calibri"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE

# Client name
client = slide.shapes.add_textbox(MARGIN_L, Emu(2286000), CONTENT_W, Emu(457200))
tf = client.text_frame
p = tf.paragraphs[0]
p.text = "[Client Name]"
p.font.name = "Calibri"
p.font.size = Pt(20)
p.font.color.rgb = DARK

# Divider
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(2926080),
                                  Emu(2286000), Emu(25400))
divider.fill.solid()
divider.fill.fore_color.rgb = BRIGHT
divider.line.fill.background()

# Subtitle
sub = slide.shapes.add_textbox(MARGIN_L, Emu(3200400), CONTENT_W, Emu(548640))
tf = sub.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ClearLaunch System  |  Step 5: Offer Development\n[Date]"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = DARK

add_footer(slide, "Offer Dev Summary")

# ========================================
# SLIDE 2: The Offer Ladder (3-tier visual)
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "THE OFFER LADDER")

# Three stacked tier cards — Core on top, Macro middle, Micro bottom
# Each with numbered badge (3 = Core, 2 = Macro, 1 = Micro)
tier_top = 990600
tier_h = 990600
tier_gap = 76200

tiers = [
    (3, "CORE OFFER (Full Engagement)",
     "[Primary revenue driver — retainer, subscription, full project scope, multi-year contract. The natural destination of the ladder.]"),
    (2, "MACRO OFFER (Entry-Point Paid)",
     "[Low-risk paid engagement — consultation, pilot, starter package, trial. Proves the relationship works and opens pathway to Core.]"),
    (1, "MICRO OFFER (Free Resource)",
     "[Free value exchange — audit, calculator, quiz, guide, sample. Demonstrates expertise and captures qualified leads.]"),
]

for i, (num, name, desc) in enumerate(tiers):
    y = tier_top + i * (tier_h + tier_gap)

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, y,
                                   CONTENT_W, Emu(tier_h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE if num != 3 else LGGREEN
    card.line.color.rgb = BORDER if num != 3 else GREEN
    card.line.width = Pt(0.5 if num != 3 else 1)

    # Left accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Emu(y + 25400),
                                     Emu(50800), Emu(tier_h - 50800))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    # Numbered badge
    add_numbered_badge(slide, int(MARGIN_L) + 190500, y + 127000, num)

    # Tier name
    name_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 660400), Emu(y + 127000),
                                         Emu(int(CONTENT_W) - 838200), Emu(254000))
    tf = name_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.name = "Calibri"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Description
    desc_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 660400), Emu(y + 431800),
                                         Emu(int(CONTENT_W) - 838200), Emu(tier_h - 533400))
    tf2 = desc_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.name = "Calibri"
    p2.font.size = Pt(10)
    p2.font.color.rgb = DARK

add_footer(slide)

# ========================================
# SLIDE 3: Micro Offer
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "MICRO OFFER (Free Entry Point)")

card_top = 1051560
card_h = 1143000
card_gap = 152400

add_content_card(slide, MARGIN_L, card_top, Emu(3931920), card_h,
                 "Format & Type",
                 "[What format will the Micro Offer take? (audit, checklist, guide, webinar, calculator, quiz, sample kit, ROI tool) Include specific title/name.]")

add_content_card(slide, Emu(4663440), card_top, Emu(3931920), card_h,
                 "Pain Point Addressed",
                 "[Which Tier 1 pain point from the ICP does this resource directly address? What common misconception or mistake does it correct?]")

add_content_card(slide, MARGIN_L, card_top + card_h + card_gap,
                 CONTENT_W, Emu(1600200),
                 "Path to Paid (How This Leads to Macro Offer)",
                 "[What actionable value does this deliver that demonstrates your unique approach? How does this resource naturally lead to a conversation about the Macro Offer?]")

add_footer(slide)

# ========================================
# SLIDE 4: Macro Offer
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "MACRO OFFER (Entry-Point Paid)")

add_content_card(slide, MARGIN_L, card_top, Emu(3931920), card_h,
                 "Engagement Type & Pricing",
                 "[What type of engagement? (consultation, assessment, pilot, starter package, diagnostic) What is the pricing strategy? (pilot pricing, introductory rate, flat fee, per-unit)]")

add_content_card(slide, Emu(4663440), card_top, Emu(3931920), card_h,
                 "Deliverable / Value Shown",
                 "[What specific deliverable or experience showcases what makes the business different? What diagnostic reveals deeper needs only this business can address?]")

add_content_card(slide, MARGIN_L, card_top + card_h + card_gap,
                 CONTENT_W, Emu(1600200),
                 "Conversion Path to Core Offer",
                 "[Typical conversion rate from Macro → Core. Natural upsell/expansion opportunities. How ROI is framed vs. cost. Common objections handled in this step.]")

add_footer(slide)

# ========================================
# SLIDE 5: Core Offer
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "CORE OFFER (Full Engagement)")

add_content_card(slide, MARGIN_L, card_top, Emu(3931920), card_h,
                 "Scope & Deliverables",
                 "[Full scope of the core offering. Specific deliverables, outcomes, or access the customer receives. What makes this the natural next step from the Macro Offer.]")

add_content_card(slide, Emu(4663440), card_top, Emu(3931920), card_h,
                 "Pricing Structure & Term",
                 "[Pricing structure (retainer, project-based, subscription, per-unit, tiered). Typical contract length or commitment period.]")

add_content_card(slide, MARGIN_L, card_top + card_h + card_gap,
                 CONTENT_W, Emu(1600200),
                 "Retention Drivers & Expansion",
                 "[What ongoing value keeps customers engaged long-term? What expansion or upsell opportunities exist within the Core Offer (additional services, higher tiers, add-ons)?]")

add_footer(slide)

# ========================================
# SLIDE 6: Creative Angles (4 columns)
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "CREATIVE ANGLES")

col_w = Emu(1905000)
col_gap = Emu(110490)

for i in range(4):
    x = int(MARGIN_L) + i * (int(col_w) + int(col_gap))
    y = card_top

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Emu(2926080))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(0.5)

    # Number badge
    add_numbered_badge(slide, x + 127000, y + 127000, i + 1)

    # Pain Point label
    lbl_box = slide.shapes.add_textbox(x + 127000, y + 482600, int(col_w) - 254000, 228600)
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Pain Point (from ICP)"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Pain Point text
    pp_box = slide.shapes.add_textbox(x + 127000, y + 711200, int(col_w) - 254000, Emu(533400))
    tf2 = pp_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "[Tier 1 pain point]"
    p2.font.name = "Calibri"
    p2.font.size = Pt(9)
    p2.font.color.rgb = DARK

    # Hook label
    hk_lbl = slide.shapes.add_textbox(x + 127000, y + 1295400, int(col_w) - 254000, 228600)
    tf = hk_lbl.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Hook Direction"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Hook text
    hk_box = slide.shapes.add_textbox(x + 127000, y + 1524000, int(col_w) - 254000, Emu(533400))
    tf2 = hk_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "[Hook using UVP language]"
    p2.font.name = "Calibri"
    p2.font.size = Pt(9)
    p2.font.color.rgb = DARK

    # CTA label
    cta_lbl = slide.shapes.add_textbox(x + 127000, y + 2108200, int(col_w) - 254000, 228600)
    tf = cta_lbl.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CTA"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # CTA text
    cta_box = slide.shapes.add_textbox(x + 127000, y + 2336800, int(col_w) - 254000, Emu(482600))
    tf2 = cta_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "[Call to action]"
    p2.font.name = "Calibri"
    p2.font.size = Pt(9)
    p2.font.color.rgb = DARK

add_footer(slide)

# ========================================
# SLIDE 7: Objection Handling (5 pairs)
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "OBJECTION HANDLING")

# Two columns: Objection | Response Framework
obj_col_w = Emu(2819400)
resp_col_w = Emu(5080000)
row_h = Emu(558800)
row_gap = Emu(25400)
row_top = 1005840

# Header row
hdr_obj = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, row_top,
                                  obj_col_w, Emu(304800))
hdr_obj.fill.solid()
hdr_obj.fill.fore_color.rgb = GREEN
hdr_obj.line.fill.background()
tf = hdr_obj.text_frame
tf.margin_left = Emu(101600)
p = tf.paragraphs[0]
p.text = "Objection"
p.font.name = "Calibri"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WHITE

hdr_resp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   int(MARGIN_L) + int(obj_col_w) + 50800,
                                   row_top, resp_col_w, Emu(304800))
hdr_resp.fill.solid()
hdr_resp.fill.fore_color.rgb = GREEN
hdr_resp.line.fill.background()
tf = hdr_resp.text_frame
tf.margin_left = Emu(101600)
p = tf.paragraphs[0]
p.text = "Response Framework (Grounded in UVP)"
p.font.name = "Calibri"
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = WHITE

# 5 objection rows
for i in range(5):
    y = row_top + 330200 + i * (int(row_h) + int(row_gap))
    fill_color = ALTROW if i % 2 == 0 else WHITE

    # Objection cell
    obj_cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, y, obj_col_w, row_h)
    obj_cell.fill.solid()
    obj_cell.fill.fore_color.rgb = fill_color
    obj_cell.line.color.rgb = BORDER
    obj_cell.line.width = Pt(0.25)
    tf = obj_cell.text_frame
    tf.margin_left = Emu(101600)
    tf.margin_top = Emu(76200)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{i+1}. [Objection {i+1}]"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = DARK

    # Response cell
    resp_cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        int(MARGIN_L) + int(obj_col_w) + 50800,
                                        y, resp_col_w, row_h)
    resp_cell.fill.solid()
    resp_cell.fill.fore_color.rgb = fill_color
    resp_cell.line.color.rgb = BORDER
    resp_cell.line.width = Pt(0.25)
    tf = resp_cell.text_frame
    tf.margin_left = Emu(101600)
    tf.margin_top = Emu(76200)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Response referencing UVP differentiator + specific proof point]"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = DARK

add_footer(slide)

# ========================================
# SLIDE 8: Offer Ladder Synthesis
# ========================================
slide = prs.slides.add_slide(blank_layout)
add_header_banner(slide, "OFFER LADDER SYNTHESIS")

# Synthesis card
card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, card_top,
                                CONTENT_W, Emu(1371600))
card1.fill.solid()
card1.fill.fore_color.rgb = LGGREEN
card1.line.color.rgb = GREEN
card1.line.width = Pt(1)

# Accent bar
accent1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, card_top + 25400,
                                  Emu(50800), Emu(1320800))
accent1.fill.solid()
accent1.fill.fore_color.rgb = GREEN
accent1.line.fill.background()

# Label
lbl1 = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), card_top + 76200,
                                  Emu(7620000), Emu(228600))
tf = lbl1.text_frame
p = tf.paragraphs[0]
p.text = "How The Ladder Connects"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GREEN

# Synthesis text
syn_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), card_top + 355600,
                                     Emu(7620000), Emu(914400))
tf = syn_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[Explain how each tier creates enough value and trust that the next step feels like an obvious decision: Micro demonstrates expertise → builds enough credibility that a prospect pays for the Macro → Macro diagnostic reveals depth of need that makes Core the natural solution.]"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = DARK

# Positioning card (links back to UVP)
pos_top = card_top + 1371600 + 152400
card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, pos_top,
                                CONTENT_W, Emu(1371600))
card2.fill.solid()
card2.fill.fore_color.rgb = WHITE
card2.line.color.rgb = BORDER
card2.line.width = Pt(0.5)

accent2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, pos_top + 25400,
                                  Emu(50800), Emu(1320800))
accent2.fill.solid()
accent2.fill.fore_color.rgb = GREEN
accent2.line.fill.background()

lbl2 = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), pos_top + 76200,
                                  Emu(7620000), Emu(228600))
tf = lbl2.text_frame
p = tf.paragraphs[0]
p.text = "Positioning Statement (From UVP)"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.bold = True
p.font.color.rgb = GREEN

pos_box = slide.shapes.add_textbox(Emu(int(MARGIN_L) + 190500), pos_top + 355600,
                                    Emu(7620000), Emu(914400))
tf = pos_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "[Pull the Positioning Statement from the UVP deliverable: For [ICP] who need [primary need], [Company] is the only [category] that [key differentiator] because [proof/reason to believe].]"
p.font.name = "Calibri"
p.font.size = Pt(11)
p.font.color.rgb = DARK

add_footer(slide)

# ========================================
# SAVE
# ========================================
prs.save("Frameworks/ClearLaunch_Offer_Dev_Summary_Deck.pptx")
print("Created: Frameworks/ClearLaunch_Offer_Dev_Summary_Deck.pptx (8 slides)")
